"""
PPTX parser — python-pptx.

For each slide:
- Collect the title and all text-frame text from every shape.
- Extract any table cells.
- Extract speaker notes (``notes_slide``).

Slide body text and notes become separate :class:`TextBlock` objects so the
chunking engine can handle them independently.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation

from docsense.parsers.base import ParsedDocument, ParserRegistry, TextBlock

logger = logging.getLogger(__name__)


class PptxParser:
    """Parse ``.pptx`` files into per-slide :class:`TextBlock` objects."""

    def parse(self, filepath: Path) -> ParsedDocument:
        """
        Extract text blocks from a PowerPoint presentation.

        Parameters
        ----------
        filepath:
            Path to the ``.pptx`` file.

        Returns
        -------
        ParsedDocument
            Slide-body blocks followed by notes blocks, in slide order.
        """
        blocks: list[TextBlock] = []
        page_count = 0

        try:
            prs = Presentation(str(filepath))
            page_count = len(prs.slides)

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_lines: list[str] = []
                slide_title: str | None = None

                # ── Extract slide title ────────────────────────────────
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()
                    slide_lines.append(f"[Slide {slide_idx}] {slide_title}")

                # ── Extract all text shapes ────────────────────────────
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and text != slide_title:
                                slide_lines.append(text)

                    if shape.has_table:
                        table_text = self._table_to_text(shape.table)
                        if table_text:
                            slide_lines.append(table_text)

                if slide_lines:
                    blocks.append(TextBlock(
                        text="\n".join(slide_lines),
                        page_number=slide_idx,
                        section_title=slide_title or f"Slide {slide_idx}",
                        block_type="paragraph",
                    ))

                # ── Speaker notes ──────────────────────────────────────
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        blocks.append(TextBlock(
                            text=notes_frame.text.strip(),
                            page_number=slide_idx,
                            section_title=(
                                f"{slide_title or f'Slide {slide_idx}'} (Notes)"
                            ),
                            block_type="note",
                        ))

        except Exception as exc:
            logger.error("PPTX parse error for %s: %s", filepath, exc)
            return ParsedDocument(
                filename=filepath.name,
                file_type=".pptx",
                error=str(exc),
            )

        return ParsedDocument(
            filename=filepath.name,
            file_type=".pptx",
            blocks=blocks,
            page_count=page_count,
        )

    @staticmethod
    def _table_to_text(table) -> str:
        """Render a PPTX table as pipe-separated rows."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)


# Register with the global parser registry
ParserRegistry.register(".pptx", PptxParser())
