"""
Document text extraction for PDF, DOCX, XLSX, PPTX.

Each extractor returns a list of chunk dicts:
    {"text": str, "page": int | None}

The chunking is character-based with overlap so no chunk crosses
a token-budget limit when embedding.
"""
from __future__ import annotations
from pathlib import Path
from typing import List, Dict, Any

from app.config import CHUNK_SIZE, CHUNK_OVERLAP


# ── Chunker ───────────────────────────────────────────────────────────────────

def _chunk(text: str) -> List[str]:
    """Split text into overlapping fixed-size character chunks."""
    text = text.strip()
    if not text:
        return []
    chunks, start = [], 0
    while start < len(text):
        end   = start + CHUNK_SIZE
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= len(text):
            break
        start = end - CHUNK_OVERLAP
    return chunks


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> List[Dict[str, Any]]:
    import fitz  # pymupdf
    results = []
    doc = fitz.open(str(path))
    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text")
        for chunk in _chunk(text):
            results.append({"text": chunk, "page": page_num})
    doc.close()
    return results


def _extract_docx(path: Path) -> List[Dict[str, Any]]:
    from docx import Document
    doc  = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": c, "page": None} for c in _chunk(text)]


def _extract_xlsx(path: Path) -> List[Dict[str, Any]]:
    import openpyxl
    results = []
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    for sheet_name in wb.sheetnames:
        ws   = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(
                str(cell) if cell is not None else "" for cell in row
            )
            if row_text.strip():
                rows.append(row_text)
        sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
        for chunk in _chunk(sheet_text):
            results.append({"text": chunk, "page": None})
    wb.close()
    return results


def _extract_pptx(path: Path) -> List[Dict[str, Any]]:
    from pptx import Presentation
    results = []
    prs = Presentation(str(path))
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        slide_text = "\n".join(texts)
        for chunk in _chunk(slide_text):
            results.append({"text": chunk, "page": slide_num})
    return results


# ── Public API ────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS: set[str] = {".pdf", ".docx", ".xlsx", ".pptx"}

_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def extract(path: Path) -> List[Dict[str, Any]]:
    """
    Extract text chunks from *path*.
    Returns list of {"text": str, "page": int | None}.
    Raises ValueError for unsupported file types.
    """
    suffix = path.suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {suffix!r}")
    return extractor(path)
